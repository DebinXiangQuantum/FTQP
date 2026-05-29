from __future__ import annotations

import json
import math
import re
from dataclasses import dataclass
from pathlib import Path
from typing import Any

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parent
SPEC_DIR = ROOT / "slide_specs"
OPT_DIR = ROOT / "optimized_specs"
FIG_DIR = ROOT / "reading_share_assets" / "figures_refined"
OUT_DIR = ROOT / "reading_share_assets"
PREVIEW_DIR = OUT_DIR / "refined_previews"
PPTX_PATH = ROOT / "popl_quantum_papers_reading_share_refined_editable.pptx"
QA_PATH = OUT_DIR / "refined_qa_report.json"

SLIDE_W = 13.333333
SLIDE_H = 7.5
PX_W = 1920
PX_H = 1080
SCALE = PX_W / SLIDE_W
FONT = "Microsoft YaHei"
MIN_FONT = 18

NAVY = "0B2F5B"
INK = "17202A"
MUTED = "5C6B7A"
PALE = "F4F8FB"
PALE2 = "EAF3FA"
LINE = "C4D4E3"
WHITE = "FFFFFF"
GREEN = "EAF8F2"
WARN = "FFF7E8"


@dataclass(frozen=True)
class Paper:
    key: str
    json_name: str
    short: str
    title: str
    authors: str
    accent: str
    fig_prefix: str


PAPERS: tuple[Paper, ...] = (
    Paper("molavi", "molavi_2026_qubit_mapping_routing.json", "QMR compiler generation", "Generating Compilers for Qubit Mapping and Routing", "Molavi et al., POPL 2026", "1F6FB8", "molavi"),
    Paper("colledan", "colledan_2025_type_based_resource_estimation.json", "Type-based resource estimation", "Flexible Type-Based Resource Estimation in Quantum Circuit Description Languages", "Colledan & Dal Lago, POPL 2025", "00878F", "colledan"),
    Paper("abdulla", "abdulla_2026_parameterized_verification.json", "Parameterized verification", "Parameterized Verification of Quantum Circuits", "Abdulla et al., POPL 2026", "7C3AED", "abdulla"),
    Paper("hirata", "hirata_2025_qurts_uncomputation.json", "Qurts uncomputation", "Qurts: Automatic Quantum Uncomputation by Affine Types with Lifetime", "Hirata & Heunen, POPL 2025", "B7791F", "hirata"),
)


FIGS: dict[str, dict[int, list[str]]] = {
    "molavi": {
        1: ["overview_pipeline"],
        2: ["case_table"],
        7: ["overview_pipeline", "amaro_nisqmr_code"],
        8: ["nisqmr_problem", "scmr_problem"],
        9: ["amaro_grammar", "amaro_semantics"],
        10: ["amaro_nisqmr_code"],
        12: ["amaro_semantics"],
        13: ["case_table", "nisqmr_results"],
        14: ["scmr_raa_results", "ilq_tiqmr_results", "ablation"],
        15: ["ablation"],
    },
    "colledan": {
        1: ["variable_input"],
        3: ["teleport_circuit", "quipper_not"],
        7: ["proto_quipper_syntax", "ra_syntax"],
        8: ["racs_bundle", "wire_routing"],
        9: ["typing_rules"],
        10: ["subtyping", "semantics"],
        11: ["let_shape", "variable_input"],
        12: ["semantics"],
        13: ["qft_example", "qura_outputs"],
        14: ["grover_example"],
        15: ["qft_example", "grover_example"],
    },
    "abdulla": {
        1: ["swta_uniform"],
        3: ["state_representations", "plain_tree_automaton"],
        7: ["state_representations", "lsta"],
        8: ["plain_tree_automaton", "swta_uniform"],
        9: ["basic_transducers", "cx_transducer"],
        10: ["verification_framework"],
        11: ["qft", "maj"],
        13: ["bv_circuit", "adder_qecc"],
        14: ["grover_circuits", "hamiltonian_table"],
        15: ["hamiltonian_table"],
    },
    "hirata": {
        1: ["pebble_game"],
        2: ["language_table"],
        3: ["pebble_game"],
        7: ["syntax_types", "typing_rules"],
        8: ["language_table", "pebble_game"],
        9: ["pure_quantum", "block_function"],
        10: ["simulation_rules", "toy_eval"],
        11: ["uncomp_semantics", "drop_trait"],
        12: ["typing_rules", "drop_trait"],
        13: ["grover_code", "silq_compare"],
        14: ["uncomp_semantics"],
        15: ["silq_compare"],
    },
}

OPT_FILES = {
    "molavi": "molavi_2026_qmr_layout.json",
    "colledan": "colledan_2025_resource_layout.json",
    "abdulla": "abdulla_2026_verification_layout.json",
    "hirata": "hirata_2025_qurts_layout.json",
}

LAYOUT_BY_INDEX = {
    1: "cover",
    2: "dense_cards",
    3: "taxonomy_matrix",
    4: "taxonomy_matrix",
    5: "risk_map",
    6: "hub_spoke",
    7: "annotated_figure",
    8: "two_figures_compare",
    9: "layered_stack",
    10: "mechanism_zoom",
    11: "before_after",
    12: "evidence_table",
    13: "annotated_figure",
    14: "metric_result",
    15: "risk_map",
}

FONT_CANDIDATES = (
    "/System/Library/Fonts/Hiragino Sans GB.ttc",
    "/System/Library/Fonts/STHeiti Medium.ttc",
    "/System/Library/Fonts/Supplemental/Arial Unicode.ttf",
)


def hex_to_rgb(hex_color: str) -> tuple[int, int, int]:
    hex_color = hex_color.strip("#")
    return tuple(int(hex_color[i : i + 2], 16) for i in (0, 2, 4))


def ppt_color(hex_color: str) -> RGBColor:
    r, g, b = hex_to_rgb(hex_color)
    return RGBColor(r, g, b)


def clean(value: Any, limit: int | None = None) -> str:
    text = re.sub(r"\s+", " ", str(value or "")).strip()
    if limit and len(text) > limit:
        return text[: limit - 1].rstrip("，。；,.;") + "…"
    return text


def load_slides(paper: Paper) -> list[dict[str, Any]]:
    data = json.loads((SPEC_DIR / paper.json_name).read_text(encoding="utf-8"))
    opt_path = OPT_DIR / OPT_FILES[paper.key]
    optimized: list[dict[str, Any]] = []
    if opt_path.exists():
        optimized = json.loads(opt_path.read_text(encoding="utf-8")).get("slides", [])
    raw = data.get("slides", data if isinstance(data, list) else [])
    slides = []
    for idx, item in enumerate(raw[:15], start=1):
        opt = optimized[idx - 1] if idx - 1 < len(optimized) else {}
        bullets = item.get("bullets") or []
        slides.append(
            {
                "title": clean(item.get("title") or opt.get("title"), 56),
                "role": clean(item.get("role")),
                "main_message": clean(item.get("main_message"), 80),
                "bullets": [clean(x, 48) for x in bullets[:5]],
                "speaker_notes": clean(item.get("speaker_notes")),
                "citation": clean(item.get("citation"), 58),
                "layout": opt.get("layout") or LAYOUT_BY_INDEX[idx],
                "density": opt.get("density", "medium"),
                "visual_directive": opt.get("visual_directive", ""),
            }
        )
    if len(slides) != 15:
        raise ValueError(f"{paper.json_name}: expected 15 slides, got {len(slides)}")
    return slides


def fig_path(paper: Paper, name: str) -> Path | None:
    path = FIG_DIR / f"{paper.fig_prefix}_{name}.png"
    return path if path.exists() else None


def slide_figs(paper: Paper, local_idx: int) -> list[Path]:
    return [p for name in FIGS.get(paper.key, {}).get(local_idx, []) if (p := fig_path(paper, name))]


def add_text(slide, x: float, y: float, w: float, h: float, text: str, size: int = 18, fill: str = INK, bold: bool = False, align: PP_ALIGN = PP_ALIGN.LEFT):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.04)
    tf.margin_right = Inches(0.04)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = 1.03
    run = p.add_run()
    run.text = text
    run.font.name = FONT
    run.font.size = Pt(max(size, MIN_FONT))
    run.font.bold = bold
    run.font.color.rgb = ppt_color(fill)
    return shape


def add_rect(slide, x: float, y: float, w: float, h: float, fill: str = PALE, outline: str = LINE, radius: bool = True):
    kind = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if radius else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    shape = slide.shapes.add_shape(kind, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ppt_color(fill)
    shape.line.color.rgb = ppt_color(outline)
    shape.line.width = Pt(1)
    return shape


def add_line(slide, x1: float, y1: float, x2: float, y2: float, fill: str = LINE, width: float = 1.0):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(min(x1, x2)), Inches(min(y1, y2)), Inches(max(abs(x2 - x1), 0.01)), Inches(max(abs(y2 - y1), 0.01)))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ppt_color(fill)
    shape.line.fill.background()
    return shape


def image_size(path: Path) -> tuple[int, int]:
    with Image.open(path) as im:
        return im.size


def fit_box(path: Path, x: float, y: float, w: float, h: float) -> tuple[float, float, float, float]:
    iw, ih = image_size(path)
    scale = min(w / iw, h / ih)
    nw, nh = iw * scale, ih * scale
    return x + (w - nw) / 2, y + (h - nh) / 2, nw, nh


def add_picture_fit(slide, path: Path, x: float, y: float, w: float, h: float, caption: str = "") -> None:
    add_rect(slide, x, y, w, h, fill=WHITE, outline="AFC9DF", radius=True)
    ix, iy, iw, ih = fit_box(path, x + 0.1, y + 0.1, w - 0.2, h - 0.42)
    slide.shapes.add_picture(str(path), Inches(ix), Inches(iy), width=Inches(iw), height=Inches(ih))
    if caption:
        add_text(slide, x + 0.12, y + h - 0.3, w - 0.24, 0.24, caption, size=18, fill=MUTED)


def add_header(slide, paper: Paper, spec: dict[str, Any], global_idx: int, total: int) -> None:
    bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, Inches(SLIDE_W), Inches(0.08))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ppt_color(paper.accent)
    bar.line.fill.background()
    add_text(slide, 0.55, 0.3, 5.0, 0.28, paper.short, size=18, fill=paper.accent, bold=True)
    add_text(slide, 0.55, 0.66, 12.15, 0.62, spec["title"], size=27, fill=INK, bold=True)
    stripe = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.55), Inches(1.38), Inches(0.06), Inches(0.45))
    stripe.fill.solid()
    stripe.fill.fore_color.rgb = ppt_color(paper.accent)
    stripe.line.fill.background()
    add_text(slide, 0.72, 1.36, 11.6, 0.5, spec["main_message"], size=19, fill=NAVY, bold=True)
    add_line(slide, 0.55, 6.82, 12.78, 6.83, fill=LINE)
    add_text(slide, 0.55, 6.92, 5.8, 0.32, paper.authors, size=18, fill=MUTED)
    add_text(slide, 11.45, 6.92, 1.3, 0.32, f"{global_idx:02d} / {total:02d}", size=18, fill=paper.accent, bold=True, align=PP_ALIGN.RIGHT)


def add_bullets(slide, paper: Paper, bullets: list[str], x: float, y: float, w: float, h: float) -> None:
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.03)
    tf.margin_right = Inches(0.03)
    for idx, item in enumerate(bullets):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = f"• {item}"
        p.font.name = FONT
        p.font.size = Pt(18)
        p.font.color.rgb = ppt_color(INK)
        p.space_after = Pt(5)
        p.line_spacing = 1.03


def add_compact_rows(slide, paper: Paper, items: list[str], x: float, y: float, w: float, row_h: float = 0.42, cols: int = 1, fill_alt: bool = False) -> None:
    items = [item for item in items if item]
    gap_x = 0.18
    gap_y = 0.12
    col_w = (w - (cols - 1) * gap_x) / cols
    for idx, item in enumerate(items):
        cx = x + (idx % cols) * (col_w + gap_x)
        cy = y + (idx // cols) * (row_h + gap_y)
        if fill_alt:
            add_rect(slide, cx, cy, col_w, row_h, fill=PALE if idx % 2 else WHITE, outline="D3E1EE", radius=True)
        else:
            add_line(slide, cx, cy + row_h, cx + col_w, cy + row_h + 0.01, fill="D8E4EF")
        add_text(slide, cx + 0.1, cy + 0.08, 0.32, 0.24, f"{idx + 1}", size=18, fill=paper.accent, bold=True)
        add_text(slide, cx + 0.48, cy + 0.08, col_w - 0.55, row_h - 0.12, item, size=18, fill=INK, bold=True)


def add_tag_row(slide, paper: Paper, items: list[str], x: float, y: float, w: float, cols: int = 4) -> None:
    items = [item for item in items if item]
    gap = 0.16
    col_w = (w - (cols - 1) * gap) / cols
    for idx, item in enumerate(items):
        cx = x + (idx % cols) * (col_w + gap)
        cy = y + (idx // cols) * 0.58
        add_rect(slide, cx, cy, col_w, 0.42, fill=PALE if idx % 2 else WHITE, outline="C8DBEB", radius=True)
        add_text(slide, cx + 0.12, cy + 0.08, col_w - 0.24, 0.22, item, size=18, fill=INK, bold=True)


def add_cover(slide, paper: Paper, spec: dict[str, Any], global_idx: int, total: int, figs: list[Path]) -> None:
    band = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, Inches(2.55), Inches(SLIDE_H))
    band.fill.solid()
    band.fill.fore_color.rgb = ppt_color(paper.accent)
    band.line.fill.background()
    add_text(slide, 0.48, 0.46, 1.8, 0.5, "POPL", size=27, fill=WHITE, bold=True)
    add_text(slide, 0.48, 1.02, 1.7, 0.38, "paper reading", size=18, fill="EAF7FF")
    add_text(slide, 3.1, 0.6, 8.8, 1.22, paper.title, size=25, fill=INK, bold=True)
    add_text(slide, 3.12, 2.05, 6.4, 0.36, paper.authors, size=18, fill=paper.accent, bold=True)
    add_line(slide, 3.12, 2.5, 11.5, 2.51, fill=LINE)
    add_text(slide, 3.12, 2.85, 8.4, 0.78, spec["main_message"], size=20, fill=NAVY, bold=True)
    add_bullets(slide, paper, spec["bullets"][:3], 3.1, 4.1, 4.9, 1.4)
    if figs:
        add_picture_fit(slide, figs[0], 8.35, 3.8, 4.0, 2.2, "paper anchor")
    add_text(slide, 0.5, 6.92, 1.8, 0.3, f"{global_idx:02d} / {total:02d}", size=18, fill=WHITE, bold=True)


def card(slide, paper: Paper, x: float, y: float, w: float, h: float, text: str, idx: int, fill: str = WHITE) -> None:
    add_rect(slide, x, y, w, h, fill=fill, outline="AFC9DF", radius=True)
    add_text(slide, x + 0.12, y + 0.08, 0.35, 0.28, str(idx), size=18, fill=paper.accent, bold=True)
    add_text(slide, x + 0.48, y + 0.08, w - 0.6, h - 0.1, text, size=18, fill=INK, bold=True)


def dense_cards(slide, paper: Paper, spec: dict[str, Any], figs: list[Path]) -> None:
    bullets = (spec["bullets"] + [""] * 5)[:5]
    if figs:
        add_picture_fit(slide, figs[0], 0.72, 2.08, 4.85, 2.55, "paper evidence")
        add_compact_rows(slide, paper, bullets, 5.9, 2.08, 6.65, row_h=0.5, cols=1, fill_alt=True)
        if len(figs) > 1:
            add_picture_fit(slide, figs[1], 0.72, 4.9, 4.85, 1.15, "second crop")
    else:
        add_compact_rows(slide, paper, bullets, 0.82, 2.2, 11.45, row_h=0.55, cols=2, fill_alt=True)


def taxonomy_matrix(slide, paper: Paper, spec: dict[str, Any], figs: list[Path]) -> None:
    items = (spec["bullets"] + [""] * 4)[:4]
    x, y, w, h = 0.75, 2.12, 7.35, 3.35
    cw, ch = (w - 0.2) / 2, (h - 0.2) / 2
    for i, b in enumerate(items):
        bx, by = x + (i % 2) * (cw + 0.2), y + (i // 2) * (ch + 0.2)
        add_rect(slide, bx, by, cw, ch, fill=WHITE if i % 2 else PALE, outline="AFC9DF", radius=True)
        add_text(slide, bx + 0.18, by + 0.16, cw - 0.3, 0.28, ["对象", "约束", "机制", "证据"][i], size=18, fill=paper.accent, bold=True)
        add_text(slide, bx + 0.18, by + 0.55, cw - 0.35, ch - 0.62, b, size=18, fill=INK, bold=True)
    if figs:
        for i, fig in enumerate(figs[:2]):
            add_picture_fit(slide, fig, 8.35, 2.12 + i * 1.78, 4.0, 1.55, fig.stem.replace(f"{paper.key}_", ""))
    else:
        add_rect(slide, 8.45, 2.2, 3.7, 3.7, fill=PALE2, outline="AFC9DF", radius=True)
        add_text(slide, 8.85, 3.42, 2.9, 0.7, "二维地图\n不是顺序流程", size=22, fill=NAVY, bold=True, align=PP_ALIGN.CENTER)


def hub_spoke(slide, paper: Paper, spec: dict[str, Any], figs: list[Path]) -> None:
    add_rect(slide, 4.2, 2.42, 4.9, 1.1, fill=paper.accent, outline=paper.accent, radius=True)
    add_text(slide, 4.5, 2.65, 4.3, 0.55, clean(spec["main_message"], 38), size=20, fill=WHITE, bold=True, align=PP_ALIGN.CENTER)
    positions = [(0.78, 2.15), (9.05, 2.15), (0.78, 4.12), (9.05, 4.12)]
    for i, b in enumerate((spec["bullets"] + [""] * 4)[:4]):
        card(slide, paper, positions[i][0], positions[i][1], 3.35, 0.78, b, i + 1, fill=WHITE if i % 2 else PALE)
        add_line(slide, 4.13 if i in (0, 2) else 9.05, positions[i][1] + 0.39, 4.2 if i in (0, 2) else 9.1, 2.98, fill="AFC9DF", width=1)
    if figs:
        add_picture_fit(slide, figs[0], 4.45, 4.08, 4.35, 1.85, "zoom")


def layered_stack(slide, paper: Paper, spec: dict[str, Any], figs: list[Path]) -> None:
    labels = (spec["bullets"] + [""] * 5)[:5]
    for i, b in enumerate(labels):
        y = 2.02 + i * 0.62
        fill = [PALE2, WHITE, PALE, WHITE, WARN][i]
        add_rect(slide, 0.78 + i * 0.35, y, 10.9 - i * 0.45, 0.44, fill=fill, outline="AFC9DF", radius=True)
        add_text(slide, 1.0 + i * 0.35, y + 0.08, 10.3 - i * 0.45, 0.24, b, size=18, fill=INK, bold=True)
    if figs:
        add_picture_fit(slide, figs[0], 9.45, 4.15, 2.85, 1.55, "rule")


def two_figures_compare(slide, paper: Paper, spec: dict[str, Any], figs: list[Path]) -> None:
    if len(figs) >= 2:
        add_picture_fit(slide, figs[0], 0.72, 2.05, 5.75, 2.75, "A")
        add_picture_fit(slide, figs[1], 6.85, 2.05, 5.75, 2.75, "B")
    elif figs:
        add_picture_fit(slide, figs[0], 0.72, 2.05, 6.3, 3.0, "evidence")
    add_tag_row(slide, paper, spec["bullets"][:4], 0.8, 5.05, 11.5, cols=4)


def annotated_figure(slide, paper: Paper, spec: dict[str, Any], figs: list[Path]) -> None:
    if figs:
        add_picture_fit(slide, figs[0], 0.72, 2.03, 7.55, 4.05, "main crop")
    if len(figs) > 1:
        add_picture_fit(slide, figs[1], 8.45, 4.42, 3.9, 1.55, "second crop")
    for i, b in enumerate((spec["bullets"] + [""] * 3)[:3]):
        add_rect(slide, 8.45, 2.08 + i * 0.62, 3.88, 0.45, fill=PALE if i % 2 else WHITE, outline="AFC9DF", radius=True)
        add_text(slide, 8.62, 2.16 + i * 0.62, 3.52, 0.24, b, size=18, fill=INK, bold=True)


def mechanism_zoom(slide, paper: Paper, spec: dict[str, Any], figs: list[Path]) -> None:
    left = spec["bullets"][:3]
    add_compact_rows(slide, paper, left, 0.72, 2.18, 4.1, row_h=0.5, cols=1, fill_alt=True)
    if figs:
        add_picture_fit(slide, figs[0], 5.15, 2.05, 7.2, 2.72, "mechanism crop")
    for i, b in enumerate(spec["bullets"][3:5]):
        add_rect(slide, 5.15 + i * 3.7, 4.95, 3.45, 0.68, fill=GREEN if i == 0 else WARN, outline="AFC9DF", radius=True)
        add_text(slide, 5.35 + i * 3.7, 5.1, 3.05, 0.28, b, size=18, fill=INK, bold=True)


def before_after(slide, paper: Paper, spec: dict[str, Any], figs: list[Path]) -> None:
    labels = spec["bullets"] + [""] * 4
    for i, title in enumerate(("Before / problem", "After / paper move")):
        x = 0.78 + i * 6.1
        add_rect(slide, x, 2.08, 5.55, 3.72, fill=WHITE if i == 0 else PALE2, outline="AFC9DF", radius=True)
        add_text(slide, x + 0.22, 2.28, 4.9, 0.32, title, size=18, fill=paper.accent, bold=True)
        add_compact_rows(slide, paper, labels[i * 2 : i * 2 + 2], x + 0.25, 2.88, 5.0, row_h=0.45, cols=1, fill_alt=False)
        if len(figs) > i:
            add_picture_fit(slide, figs[i], x + 0.4, 4.25, 4.75, 1.28, f"crop {i+1}")


def evidence_table(slide, paper: Paper, spec: dict[str, Any], figs: list[Path]) -> None:
    headers = ["证据点", "为什么重要"]
    add_rect(slide, 0.78, 2.08, 7.35, 3.35, fill=WHITE, outline="AFC9DF", radius=True)
    add_text(slide, 1.0, 2.28, 2.3, 0.28, headers[0], size=18, fill=paper.accent, bold=True)
    add_text(slide, 4.1, 2.28, 2.3, 0.28, headers[1], size=18, fill=paper.accent, bold=True)
    for i, b in enumerate((spec["bullets"] + [""] * 4)[:4]):
        y = 2.78 + i * 0.58
        add_line(slide, 1.0, y - 0.08, 7.85, y - 0.07, fill="E1EAF2")
        add_text(slide, 1.0, y, 2.75, 0.32, clean(b, 22), size=18, fill=INK, bold=True)
        add_text(slide, 4.1, y, 3.55, 0.32, clean(b, 32), size=18, fill=MUTED)
    if figs:
        add_picture_fit(slide, figs[0], 8.45, 2.12, 3.9, 3.1, "formal object")


def metric_result(slide, paper: Paper, spec: dict[str, Any], figs: list[Path]) -> None:
    if figs:
        if len(figs) >= 3:
            for i, fig in enumerate(figs[:3]):
                add_picture_fit(slide, fig, 0.72 + i * 4.0, 2.05, 3.65, 2.25, fig.stem.split("_")[-1])
        elif len(figs) == 2:
            add_picture_fit(slide, figs[0], 0.72, 2.05, 5.7, 2.4, "result A")
            add_picture_fit(slide, figs[1], 6.8, 2.05, 5.7, 2.4, "result B")
        else:
            add_picture_fit(slide, figs[0], 0.72, 2.05, 6.0, 2.8, "result")
    add_tag_row(slide, paper, spec["bullets"][:4], 0.85, 5.02, 11.4, cols=4)


def risk_map(slide, paper: Paper, spec: dict[str, Any], figs: list[Path]) -> None:
    items = (spec["bullets"] + [""] * 4)[:4]
    labels = ["已解决", "仍受限", "启发", "追问"]
    fills = [GREEN, WARN, PALE2, WHITE]
    for i, b in enumerate(items):
        x = 0.78 + (i % 2) * 6.1
        y = 2.12 + (i // 2) * 1.82
        add_rect(slide, x, y, 5.55, 1.04, fill=fills[i], outline="AFC9DF", radius=True)
        add_text(slide, x + 0.22, y + 0.16, 1.2, 0.3, labels[i], size=18, fill=paper.accent, bold=True)
        add_text(slide, x + 1.35, y + 0.18, 3.85, 0.46, b, size=18, fill=INK, bold=True)
    if figs:
        add_picture_fit(slide, figs[0], 9.3, 5.4, 2.8, 0.8, "evidence")


LAYOUT_FUNCS = {
    "dense_cards": dense_cards,
    "taxonomy_matrix": taxonomy_matrix,
    "hub_spoke": hub_spoke,
    "layered_stack": layered_stack,
    "two_figures_compare": two_figures_compare,
    "annotated_figure": annotated_figure,
    "mechanism_zoom": mechanism_zoom,
    "before_after": before_after,
    "evidence_table": evidence_table,
    "metric_result": metric_result,
    "risk_map": risk_map,
}


def build_pptx(records: list[tuple[Paper, int, dict[str, Any], list[Path]]]) -> None:
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    blank = prs.slide_layouts[6]
    total = len(records)
    for global_idx, (paper, local_idx, spec, figs) in enumerate(records, start=1):
        slide = prs.slides.add_slide(blank)
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = ppt_color(WHITE)
        if spec["layout"] == "cover":
            add_cover(slide, paper, spec, global_idx, total, figs)
        else:
            add_header(slide, paper, spec, global_idx, total)
            LAYOUT_FUNCS[spec["layout"]](slide, paper, spec, figs)
            if spec["citation"]:
                add_text(slide, 0.7, 6.46, 11.6, 0.26, spec["citation"], size=18, fill=MUTED)
    prs.save(PPTX_PATH)


def pil_font(size: int, bold: bool = False) -> ImageFont.FreeTypeFont:
    path = next((p for p in FONT_CANDIDATES if Path(p).exists()), None)
    if not path:
        return ImageFont.load_default()
    try:
        return ImageFont.truetype(path, size=size, index=1 if bold and path.endswith(".ttc") else 0)
    except TypeError:
        return ImageFont.truetype(path, size=size)


def px(v: float) -> int:
    return int(round(v * SCALE))


def wrap(draw: ImageDraw.ImageDraw, text: str, fnt: ImageFont.FreeTypeFont, max_w: int) -> list[str]:
    out: list[str] = []
    cur = ""
    for ch in text:
        cand = cur + ch
        if draw.textbbox((0, 0), cand, font=fnt)[2] <= max_w or not cur:
            cur = cand
        else:
            out.append(cur)
            cur = ch
    if cur:
        out.append(cur)
    return out


def ptext(draw: ImageDraw.ImageDraw, xy: tuple[int, int], text: str, max_w: int, size: int, fill: str = INK, bold: bool = False, max_lines: int = 3) -> None:
    fnt = pil_font(max(size, 18), bold=bold)
    y = xy[1]
    lines = wrap(draw, text, fnt, max_w)[:max_lines]
    for line in lines:
        draw.text((xy[0], y), line, font=fnt, fill=hex_to_rgb(fill))
        y += int(size * 1.28)


def prect(draw: ImageDraw.ImageDraw, box: tuple[float, float, float, float], fill: str, outline: str = LINE, r: int = 16) -> None:
    draw.rounded_rectangle(tuple(px(v) for v in box), radius=r, fill=hex_to_rgb(fill), outline=hex_to_rgb(outline), width=2)


def ppaste(img: Image.Image, path: Path, box: tuple[float, float, float, float]) -> None:
    x1, y1, x2, y2 = [px(v) for v in box]
    with Image.open(path).convert("RGB") as src:
        src.thumbnail((x2 - x1, y2 - y1), Image.Resampling.LANCZOS)
        bg = Image.new("RGB", (x2 - x1, y2 - y1), "white")
        bg.paste(src, ((bg.width - src.width) // 2, (bg.height - src.height) // 2))
        img.paste(bg, (x1, y1))


def preview_slide(paper: Paper, local_idx: int, spec: dict[str, Any], figs: list[Path], global_idx: int, total: int) -> Image.Image:
    img = Image.new("RGB", (PX_W, PX_H), "white")
    draw = ImageDraw.Draw(img)
    if spec["layout"] == "cover":
        draw.rectangle((0, 0, px(2.55), PX_H), fill=hex_to_rgb(paper.accent))
        ptext(draw, (px(0.48), px(0.48)), "POPL", px(1.8), 38, WHITE, True, 1)
        ptext(draw, (px(0.48), px(1.05)), "paper reading", px(1.8), 22, "EAF7FF", False, 1)
        ptext(draw, (px(3.1), px(0.65)), paper.title, px(8.8), 34, INK, True, 3)
        ptext(draw, (px(3.12), px(2.12)), paper.authors, px(6.4), 22, paper.accent, True, 1)
        ptext(draw, (px(3.12), px(2.9)), spec["main_message"], px(8.4), 28, NAVY, True, 2)
        for i, b in enumerate(spec["bullets"][:3]):
            ptext(draw, (px(3.15), px(4.1 + i * 0.42)), f"• {b}", px(7.8), 22, INK, False, 1)
        if figs:
            prect(draw, (8.35, 3.8, 12.35, 6.0), WHITE, "AFC9DF")
            ppaste(img, figs[0], (8.45, 3.9, 12.25, 5.75))
        return img
    draw.rectangle((0, 0, PX_W, px(0.08)), fill=hex_to_rgb(paper.accent))
    ptext(draw, (px(0.55), px(0.3)), paper.short, px(5), 22, paper.accent, True, 1)
    ptext(draw, (px(0.55), px(0.68)), spec["title"], px(12), 36, INK, True, 2)
    draw.rectangle((px(0.55), px(1.38), px(0.61), px(1.83)), fill=hex_to_rgb(paper.accent))
    ptext(draw, (px(0.72), px(1.37)), spec["main_message"], px(11.6), 26, NAVY, True, 2)
    layout = spec["layout"]
    bullets = spec["bullets"]
    if layout in ("annotated_figure", "mechanism_zoom") and figs:
        prect(draw, (0.72, 2.03, 8.25, 6.05), WHITE, "AFC9DF")
        ppaste(img, figs[0], (0.86, 2.17, 8.08, 5.72))
        for i, b in enumerate(bullets[:3]):
            prect(draw, (8.45, 2.08 + i * 0.62, 12.35, 2.53 + i * 0.62), PALE if i % 2 else WHITE, "AFC9DF")
            ptext(draw, (px(8.62), px(2.15 + i * 0.62)), b, px(3.5), 22, INK, True, 1)
        if len(figs) > 1:
            prect(draw, (8.45, 4.42, 12.35, 5.98), WHITE, "AFC9DF")
            ppaste(img, figs[1], (8.55, 4.52, 12.25, 5.75))
    elif layout == "taxonomy_matrix":
        labels = ["对象", "约束", "机制", "证据"]
        for i, b in enumerate((bullets + [""] * 4)[:4]):
            x = 0.75 + (i % 2) * 3.78
            y = 2.12 + (i // 2) * 1.68
            prect(draw, (x, y, x + 3.55, y + 1.45), WHITE if i % 2 else PALE, "AFC9DF")
            ptext(draw, (px(x + 0.18), px(y + 0.16)), labels[i], px(1.1), 22, paper.accent, True, 1)
            ptext(draw, (px(x + 0.18), px(y + 0.58)), b, px(3.1), 21, INK, True, 2)
        for i, fig in enumerate(figs[:2]):
            prect(draw, (8.35, 2.12 + i * 1.78, 12.35, 3.67 + i * 1.78), WHITE, "AFC9DF")
            ppaste(img, fig, (8.45, 2.22 + i * 1.78, 12.25, 3.42 + i * 1.78))
    elif layout == "hub_spoke":
        prect(draw, (4.2, 2.42, 9.1, 3.52), paper.accent, paper.accent)
        ptext(draw, (px(4.5), px(2.68)), clean(spec["main_message"], 36), px(4.3), 25, WHITE, True, 2)
        positions = [(0.78, 2.15), (9.05, 2.15), (0.78, 4.12), (9.05, 4.12)]
        for i, b in enumerate((bullets + [""] * 4)[:4]):
            x, y = positions[i]
            prect(draw, (x, y, x + 3.35, y + 0.78), WHITE if i % 2 else PALE, "AFC9DF")
            ptext(draw, (px(x + 0.16), px(y + 0.16)), f"{i+1}  {b}", px(3.0), 21, INK, True, 1)
        if figs:
            prect(draw, (4.45, 4.08, 8.8, 5.93), WHITE, "AFC9DF")
            ppaste(img, figs[0], (4.55, 4.18, 8.7, 5.72))
    elif layout == "layered_stack":
        for i, b in enumerate((bullets + [""] * 5)[:5]):
            x = 0.78 + i * 0.35
            y = 2.02 + i * 0.62
            prect(draw, (x, y, 11.68 - i * 0.1, y + 0.44), [PALE2, WHITE, PALE, WHITE, WARN][i], "AFC9DF")
            ptext(draw, (px(x + 0.22), px(y + 0.08)), b, px(9.8 - i * 0.2), 21, INK, True, 1)
        if figs:
            prect(draw, (9.45, 4.15, 12.3, 5.7), WHITE, "AFC9DF")
            ppaste(img, figs[0], (9.55, 4.25, 12.2, 5.48))
    elif layout in ("dense_cards", "evidence_table", "risk_map", "before_after"):
        for i, b in enumerate(bullets[:5]):
            x = 0.78 + (i % 2) * 5.9
            y = 2.1 + (i // 2) * 0.72
            prect(draw, (x, y, x + 5.45, y + 0.52), PALE if i % 2 else WHITE, "AFC9DF")
            ptext(draw, (px(x + 0.18), px(y + 0.12)), f"{i+1}  {b}", px(5.0), 21, INK, True, 1)
        for i, fig in enumerate(figs[:2]):
            prect(draw, (0.78 + i * 5.9, 4.55, 6.15 + i * 5.9, 6.05), WHITE, "AFC9DF")
            ppaste(img, fig, (0.88 + i * 5.9, 4.65, 6.05 + i * 5.9, 5.82))
    elif layout in ("two_figures_compare", "metric_result") and figs:
        for i, fig in enumerate(figs[:3]):
            x = 0.72 + i * 4.0 if len(figs) >= 3 else 0.72 + i * 6.1
            w = 3.65 if len(figs) >= 3 else 5.75
            prect(draw, (x, 2.05, x + w, 4.75), WHITE, "AFC9DF")
            ppaste(img, fig, (x + 0.1, 2.15, x + w - 0.1, 4.48))
        for i, b in enumerate(bullets[:4]):
            prect(draw, (0.85 + i * 3, 5.02, 3.5 + i * 3, 5.48), PALE if i % 2 else WHITE, "AFC9DF")
            ptext(draw, (px(1.0 + i * 3), px(5.12)), b, px(2.35), 20, INK, True, 1)
    else:
        for i, b in enumerate(bullets[:5]):
            x = 0.78 + (i % 2) * 5.8
            y = 2.1 + (i // 2) * 0.72
            prect(draw, (x, y, x + 5.35, y + 0.52), PALE if i % 2 else WHITE, "AFC9DF")
            ptext(draw, (px(x + 0.2), px(y + 0.12)), b, px(4.9), 22, INK, True, 1)
        if figs:
            prect(draw, (8.35, 4.2, 12.2, 6.0), WHITE, "AFC9DF")
            ppaste(img, figs[0], (8.45, 4.3, 12.1, 5.75))
    draw.line((px(0.55), px(6.82), px(12.78), px(6.82)), fill=hex_to_rgb(LINE), width=2)
    ptext(draw, (px(0.55), px(6.92)), paper.authors, px(5.8), 22, MUTED, False, 1)
    ptext(draw, (px(11.55), px(6.92)), f"{global_idx:02d} / {total:02d}", px(1.2), 22, paper.accent, True, 1)
    return img


def build_previews(records: list[tuple[Paper, int, dict[str, Any], list[Path]]]) -> Path:
    PREVIEW_DIR.mkdir(parents=True, exist_ok=True)
    paths: list[Path] = []
    total = len(records)
    for global_idx, (paper, local_idx, spec, figs) in enumerate(records, start=1):
        img = preview_slide(paper, local_idx, spec, figs, global_idx, total)
        out = PREVIEW_DIR / f"slide_{global_idx:02d}.png"
        img.save(out, "PNG", optimize=True)
        paths.append(out)
    thumb_w, thumb_h = 384, 216
    cols = 5
    rows = math.ceil(len(paths) / cols)
    sheet = Image.new("RGB", (cols * thumb_w, rows * (thumb_h + 34)), hex_to_rgb("EEF3F7"))
    draw = ImageDraw.Draw(sheet)
    for idx, path in enumerate(paths):
        with Image.open(path).convert("RGB") as src:
            src.thumbnail((thumb_w, thumb_h), Image.Resampling.LANCZOS)
            x = (idx % cols) * thumb_w
            y = (idx // cols) * (thumb_h + 34)
            sheet.paste(src, (x, y))
            draw.text((x + 10, y + thumb_h + 6), f"{idx + 1:02d}", font=pil_font(22, True), fill=hex_to_rgb(NAVY))
    contact = PREVIEW_DIR / "contact_sheet.png"
    sheet.save(contact, "PNG", optimize=True)
    return contact


def main() -> None:
    OUT_DIR.mkdir(parents=True, exist_ok=True)
    records: list[tuple[Paper, int, dict[str, Any], list[Path]]] = []
    for paper in PAPERS:
        for local_idx, spec in enumerate(load_slides(paper), start=1):
            records.append((paper, local_idx, spec, slide_figs(paper, local_idx)))
    build_pptx(records)
    contact = build_previews(records)
    QA_PATH.write_text(
        json.dumps(
            {
                "pptx": str(PPTX_PATH),
                "slide_count": len(records),
                "preview_contact_sheet": str(contact),
                "preview_dir": str(PREVIEW_DIR),
                "figure_source": str(FIG_DIR),
                "min_font_pt": MIN_FONT,
            },
            ensure_ascii=False,
            indent=2,
        ),
        encoding="utf-8",
    )
    print(f"Wrote refined editable PPTX: {PPTX_PATH}")
    print(f"Wrote previews: {PREVIEW_DIR}")
    print(f"Wrote contact sheet: {contact}")


if __name__ == "__main__":
    main()
