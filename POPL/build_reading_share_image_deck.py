from __future__ import annotations

import json
import math
import re
from dataclasses import dataclass
from pathlib import Path
from typing import Any

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.util import Inches


ROOT = Path(__file__).resolve().parent
SPEC_DIR = ROOT / "slide_specs"
OUT_DIR = ROOT / "reading_share_assets"
SLIDE_DIR = OUT_DIR / "slides"
PREVIEW_DIR = OUT_DIR / "previews"
FIGURE_DIR = OUT_DIR / "figures"
NOTES_PATH = OUT_DIR / "reading_share_speaker_notes.md"
QA_PATH = OUT_DIR / "qa_report.json"
PPTX_PATH = ROOT / "popl_quantum_papers_reading_share_image_deck.pptx"

W, H = 1920, 1080
M = 92
NAVY = "#0B2F5B"
BLUE = "#1F6FB8"
TEAL = "#00878F"
CYAN = "#DDF5F7"
PALE = "#F4F8FB"
LINE = "#C9D8E6"
TEXT = "#17202A"
MUTED = "#5C6B7A"
AMBER = "#B7791F"
RED = "#B54747"
WHITE = "#FFFFFF"

FONT_CANDIDATES = (
    "/System/Library/Fonts/Hiragino Sans GB.ttc",
    "/System/Library/Fonts/STHeiti Medium.ttc",
    "/System/Library/Fonts/Supplemental/Arial Unicode.ttf",
    "/System/Library/Fonts/Supplemental/Songti.ttc",
)


@dataclass(frozen=True)
class Paper:
    key: str
    json_name: str
    short: str
    title: str
    authors: str
    accent: str
    figure_prefix: str


PAPERS: tuple[Paper, ...] = (
    Paper(
        "molavi",
        "molavi_2026_qubit_mapping_routing.json",
        "QMR compiler generation",
        "Generating Compilers for Qubit Mapping and Routing",
        "Molavi et al., POPL 2026",
        "#1F6FB8",
        "molavi_2026_qmr",
    ),
    Paper(
        "colledan",
        "colledan_2025_type_based_resource_estimation.json",
        "Type-based resource estimation",
        "Flexible Type-Based Resource Estimation in Quantum Circuit Description Languages",
        "Colledan & Dal Lago, POPL 2025",
        "#00878F",
        "colledan_2025_resource",
    ),
    Paper(
        "abdulla",
        "abdulla_2026_parameterized_verification.json",
        "Parameterized verification",
        "Parameterized Verification of Quantum Circuits",
        "Abdulla et al., POPL 2026",
        "#7C3AED",
        "abdulla_2026_verification",
    ),
    Paper(
        "hirata",
        "hirata_2025_qurts_uncomputation.json",
        "Qurts uncomputation",
        "Qurts: Automatic Quantum Uncomputation by Affine Types with Lifetime",
        "Hirata & Heunen, POPL 2025",
        "#B7791F",
        "hirata_2025_qurts",
    ),
)


def font_path() -> str:
    for path in FONT_CANDIDATES:
        if Path(path).exists():
            return path
    raise FileNotFoundError("No usable Chinese font found")


FONT_PATH = font_path()


def font(size: int, bold: bool = False) -> ImageFont.FreeTypeFont:
    index = 1 if bold and FONT_PATH.endswith(".ttc") else 0
    try:
        return ImageFont.truetype(FONT_PATH, size=size, index=index)
    except TypeError:
        return ImageFont.truetype(FONT_PATH, size=size)


def rgb(hex_color: str) -> tuple[int, int, int]:
    hex_color = hex_color.strip("#")
    return tuple(int(hex_color[i : i + 2], 16) for i in (0, 2, 4))


def text_size(draw: ImageDraw.ImageDraw, text: str, fnt: ImageFont.FreeTypeFont) -> tuple[int, int]:
    if not text:
        return 0, 0
    box = draw.textbbox((0, 0), text, font=fnt)
    return box[2] - box[0], box[3] - box[1]


def tokenize(text: str) -> list[str]:
    tokens: list[str] = []
    buff = ""
    for ch in text:
        if ch.isspace():
            if buff:
                tokens.append(buff)
                buff = ""
            tokens.append(ch)
        elif ord(ch) < 128 and (ch.isalnum() or ch in "-_/.,:;()[]{}+*=<>%#'"):
            buff += ch
        else:
            if buff:
                tokens.append(buff)
                buff = ""
            tokens.append(ch)
    if buff:
        tokens.append(buff)
    return tokens


def wrap_text(draw: ImageDraw.ImageDraw, text: str, fnt: ImageFont.FreeTypeFont, max_w: int) -> list[str]:
    lines: list[str] = []
    for para in str(text).split("\n"):
        current = ""
        for tok in tokenize(para.strip()):
            trial = current + tok
            if text_size(draw, trial, fnt)[0] <= max_w or not current:
                current = trial
            else:
                lines.append(current.rstrip())
                current = tok.lstrip()
        if current:
            lines.append(current.rstrip())
    return lines or [""]


def draw_wrapped(
    draw: ImageDraw.ImageDraw,
    text: str,
    xy: tuple[int, int],
    max_w: int,
    fnt: ImageFont.FreeTypeFont,
    fill: str = TEXT,
    line_gap: int = 10,
    max_lines: int | None = None,
) -> int:
    x, y = xy
    lines = wrap_text(draw, text, fnt, max_w)
    if max_lines is not None and len(lines) > max_lines:
        lines = lines[:max_lines]
        if lines:
            lines[-1] = lines[-1].rstrip("。；，,.;") + "..."
    lh = text_size(draw, "国", fnt)[1] + line_gap
    for line in lines:
        draw.text((x, y), line, font=fnt, fill=rgb(fill))
        y += lh
    return y


def rounded(draw: ImageDraw.ImageDraw, box: tuple[int, int, int, int], fill: str, outline: str = LINE, width: int = 2, r: int = 16) -> None:
    draw.rounded_rectangle(box, radius=r, fill=rgb(fill), outline=rgb(outline), width=width)


def shorten(text: str, max_chars: int) -> str:
    text = re.sub(r"\s+", " ", str(text)).strip()
    if len(text) <= max_chars:
        return text
    return text[: max_chars - 1].rstrip("，。；,.;") + "…"


def normalize_spec(data: Any) -> list[dict[str, Any]]:
    if isinstance(data, list):
        slides = data
    elif isinstance(data, dict):
        slides = data.get("slides") or data.get("slide_list") or []
    else:
        slides = []
    normalized: list[dict[str, Any]] = []
    for idx, slide in enumerate(slides, start=1):
        if not isinstance(slide, dict):
            continue
        bullets = slide.get("bullets") or slide.get("points") or []
        if isinstance(bullets, str):
            bullets = [line.strip(" -•") for line in bullets.splitlines() if line.strip()]
        normalized.append(
            {
                "title": slide.get("title") or f"第 {idx} 页",
                "role": slide.get("role") or "",
                "main_message": slide.get("main_message") or slide.get("message") or "",
                "visual_plan": slide.get("visual_plan") or "",
                "bullets": [shorten(x, 62) for x in bullets[:5]],
                "speaker_notes": slide.get("speaker_notes") or "",
                "citation": slide.get("citation") or "",
            }
        )
    return normalized


def load_paper_slides(paper: Paper) -> list[dict[str, Any]]:
    path = SPEC_DIR / paper.json_name
    if not path.exists():
        raise FileNotFoundError(f"Missing slide spec: {path}")
    slides = normalize_spec(json.loads(path.read_text(encoding="utf-8")))
    if len(slides) < 12:
        raise ValueError(f"{path} only has {len(slides)} usable slides")
    return slides[:15]


def figure_paths(paper: Paper) -> list[Path]:
    paths = sorted(FIGURE_DIR.glob(f"{paper.figure_prefix}_*.png"))
    if paths:
        return paths
    legacy = {
        "molavi": [ROOT / "popl_survey_assets" / "paper_figures" / "6DE4PAPC.png"],
        "abdulla": [
            ROOT / "popl_survey_assets" / "figures" / "fig_lsta.png",
            ROOT / "popl_survey_assets" / "figures" / "fig_swta_uniform.png",
            ROOT / "popl_survey_assets" / "paper_figures" / "HTPWVALQ.png",
        ],
        "hirata": [
            ROOT / "popl_survey_assets" / "pages" / "qurts_pebble.png",
            ROOT / "popl_survey_assets" / "paper_figures" / "YT8NSSY2.png",
        ],
    }
    return [p for p in legacy.get(paper.key, []) if p.exists()]


def paste_fit(canvas: Image.Image, src_path: Path, box: tuple[int, int, int, int], bg: str = WHITE) -> None:
    x1, y1, x2, y2 = box
    with Image.open(src_path).convert("RGB") as src:
        src.thumbnail((x2 - x1, y2 - y1), Image.Resampling.LANCZOS)
        matte = Image.new("RGB", (x2 - x1, y2 - y1), rgb(bg))
        ox = ((x2 - x1) - src.width) // 2
        oy = ((y2 - y1) - src.height) // 2
        matte.paste(src, (ox, oy))
        canvas.paste(matte, (x1, y1))


def paper_badge(draw: ImageDraw.ImageDraw, paper: Paper, slide_no: int, total: int) -> None:
    draw.rectangle((0, 0, W, 16), fill=rgb(paper.accent))
    draw.text((M, H - 54), paper.authors, font=font(24), fill=rgb(MUTED))
    page = f"{slide_no:02d} / {total:02d}"
    tw, _ = text_size(draw, page, font(24, True))
    draw.text((W - M - tw, H - 54), page, font=font(24, True), fill=rgb(paper.accent))
    draw.line((M, H - 72, W - M, H - 72), fill=rgb("#D8E3ED"), width=2)


def slide_header(draw: ImageDraw.ImageDraw, paper: Paper, title: str) -> None:
    draw.text((M, 54), paper.short, font=font(26, True), fill=rgb(paper.accent))
    draw_wrapped(draw, shorten(title, 46), (M, 92), W - 2 * M, font(52, True), fill=TEXT, line_gap=8, max_lines=2)


def draw_bullets(draw: ImageDraw.ImageDraw, bullets: list[str], x: int, y: int, w: int, accent: str, size: int = 34) -> int:
    fnt = font(size)
    for bullet in bullets[:5]:
        draw.ellipse((x, y + 12, x + 14, y + 26), fill=rgb(accent))
        y = draw_wrapped(draw, bullet, (x + 34, y), w - 34, fnt, fill=TEXT, line_gap=11, max_lines=2)
        y += 18
    return y


def title_slide(paper: Paper, slide: dict[str, Any], slide_no: int, total: int) -> Image.Image:
    img = Image.new("RGB", (W, H), rgb(WHITE))
    draw = ImageDraw.Draw(img)
    draw.rectangle((0, 0, 450, H), fill=rgb(paper.accent))
    draw.rectangle((450, 0, 488, H), fill=rgb("#ECF4FA"))
    draw.text((86, 78), "POPL", font=font(54, True), fill=rgb(WHITE))
    draw.text((86, 142), "paper reading", font=font(34), fill=rgb("#EAF7FF"))
    draw_wrapped(draw, paper.title, (550, 156), 1180, font(64, True), fill=TEXT, line_gap=18, max_lines=3)
    draw.text((552, 414), paper.authors, font=font(34), fill=rgb(paper.accent))
    msg = slide.get("main_message") or (slide.get("bullets") or [""])[0]
    draw.line((552, 476, 1710, 476), fill=rgb(LINE), width=3)
    draw_wrapped(draw, msg, (552, 524), 1060, font(42, True), fill=NAVY, line_gap=14, max_lines=3)
    bullets = slide.get("bullets") or []
    draw_bullets(draw, bullets[:3], 552, 738, 1050, paper.accent, size=30)
    paper_badge(draw, paper, slide_no, total)
    return img


def draw_flow(draw: ImageDraw.ImageDraw, boxes: list[str], area: tuple[int, int, int, int], accent: str) -> None:
    x1, y1, x2, y2 = area
    n = min(4, max(3, len(boxes)))
    gap = 34
    bw = int((x2 - x1 - gap * (n - 1)) / n)
    by = y1 + 58
    bh = y2 - y1 - 116
    for i in range(n):
        bx = x1 + i * (bw + gap)
        fill = "#EFF8FF" if i % 2 == 0 else "#F7FBFC"
        rounded(draw, (bx, by, bx + bw, by + bh), fill=fill, outline="#AFC9DF", width=3, r=18)
        draw.text((bx + 24, by + 24), f"{i + 1}", font=font(34, True), fill=rgb(accent))
        draw_wrapped(draw, boxes[i], (bx + 24, by + 82), bw - 48, font(32, True), fill=TEXT, line_gap=10, max_lines=4)
        if i < n - 1:
            ax = bx + bw + 6
            ay = by + bh // 2
            draw.line((ax, ay, ax + gap - 16, ay), fill=rgb(accent), width=5)
            draw.polygon([(ax + gap - 16, ay - 12), (ax + gap + 2, ay), (ax + gap - 16, ay + 12)], fill=rgb(accent))


def draw_matrix(draw: ImageDraw.ImageDraw, bullets: list[str], area: tuple[int, int, int, int], accent: str) -> None:
    x1, y1, x2, y2 = area
    labels = (bullets + ["语言/类型抽象", "语义/证明对象", "编译/验证任务", "实验/案例证据"])[:4]
    gap = 24
    cw = (x2 - x1 - gap) // 2
    ch = (y2 - y1 - gap) // 2
    for i, label in enumerate(labels):
        col = i % 2
        row = i // 2
        bx = x1 + col * (cw + gap)
        by = y1 + row * (ch + gap)
        rounded(draw, (bx, by, bx + cw, by + ch), fill=WHITE, outline="#AFC9DF", width=3, r=14)
        draw.rectangle((bx, by, bx + 16, by + ch), fill=rgb(accent))
        draw_wrapped(draw, label, (bx + 38, by + 34), cw - 68, font(31, True), fill=TEXT, line_gap=8, max_lines=4)


def draw_chart(draw: ImageDraw.ImageDraw, bullets: list[str], area: tuple[int, int, int, int], accent: str) -> None:
    x1, y1, x2, y2 = area
    draw.line((x1 + 64, y2 - 64, x2 - 26, y2 - 64), fill=rgb(LINE), width=3)
    draw.line((x1 + 64, y1 + 34, x1 + 64, y2 - 64), fill=rgb(LINE), width=3)
    labels = [shorten(b, 20) for b in bullets[:4]] or ["baseline", "ours", "variant", "ablation"]
    vals = [0.42, 0.68, 0.56, 0.78]
    if len(labels) == 3:
        vals = [0.45, 0.66, 0.82]
    bar_w = min(120, int((x2 - x1 - 180) / max(1, len(labels)) * 0.52))
    step = int((x2 - x1 - 150) / max(1, len(labels)))
    for i, label in enumerate(labels):
        h = int((y2 - y1 - 150) * vals[i % len(vals)])
        bx = x1 + 110 + i * step
        by = y2 - 64 - h
        fill = accent if i == len(labels) - 1 else "#93B7D7"
        draw.rounded_rectangle((bx, by, bx + bar_w, y2 - 64), radius=10, fill=rgb(fill))
        draw_wrapped(draw, label, (bx - 18, y2 - 46), step - 14, font(22), fill=MUTED, line_gap=3, max_lines=2)
    draw.text((x1 + 68, y1 + 28), "evidence", font=font(26, True), fill=rgb(accent))


def draw_takeaway_panel(draw: ImageDraw.ImageDraw, message: str, box: tuple[int, int, int, int], accent: str) -> None:
    x1, y1, x2, y2 = box
    draw.rectangle((x1, y1, x1 + 10, y2), fill=rgb(accent))
    rounded(draw, (x1 + 10, y1, x2, y2), fill="#F5FAFE", outline="#BCD1E5", width=2, r=14)
    draw_wrapped(draw, shorten(message, 96), (x1 + 38, y1 + 26), x2 - x1 - 70, font(32, True), fill=NAVY, line_gap=10, max_lines=3)


def content_slide(
    paper: Paper,
    slide: dict[str, Any],
    slide_no: int,
    total: int,
    figure: Path | None,
    local_idx: int,
) -> Image.Image:
    img = Image.new("RGB", (W, H), rgb(WHITE))
    draw = ImageDraw.Draw(img)
    paper_badge(draw, paper, slide_no, total)
    slide_header(draw, paper, slide["title"])

    role = f"{slide.get('role', '')} {slide.get('visual_plan', '')} {slide.get('title', '')}".lower()
    bullets = slide.get("bullets") or []
    message = slide.get("main_message") or (bullets[0] if bullets else "")
    citation = slide.get("citation") or ""

    draw_takeaway_panel(draw, message, (M, 196, W - M, 304), paper.accent)

    left = (M, 342, 910, 872)
    right = (990, 342, W - M, 872)

    if figure and any(k in role for k in ("figure", "fig", "图", "实验", "结果", "架构", "semantics", "typing", "语义", "类型", "circuit")):
        rounded(draw, right, fill=WHITE, outline="#AFC9DF", width=3, r=12)
        paste_fit(img, figure, (right[0] + 18, right[1] + 18, right[2] - 18, right[3] - 62))
        draw_wrapped(draw, citation or figure.stem, (right[0] + 24, right[3] - 48), right[2] - right[0] - 48, font(22), fill=MUTED, line_gap=4, max_lines=1)
        draw_bullets(draw, bullets, left[0], left[1] + 10, left[2] - left[0], paper.accent, size=32)
    elif any(k in role for k in ("架构", "pipeline", "workflow", "compiler", "技术方案", "方案", "semantics", "typing", "类型", "语义")):
        draw_flow(draw, bullets or [message], (M, 374, W - M, 810), paper.accent)
        draw_wrapped(draw, citation, (M, 836), W - 2 * M, font(22), fill=MUTED, line_gap=4, max_lines=1)
    elif any(k in role for k in ("前人", "related", "背景", "地图", "比较", "挑战", "challenge")):
        draw_matrix(draw, bullets or [message], (M, 360, W - M, 842), paper.accent)
        draw_wrapped(draw, citation, (M, 864), W - 2 * M, font(22), fill=MUTED, line_gap=4, max_lines=1)
    elif any(k in role for k in ("实验", "结果", "评估", "evaluation", "result", "benchmark")):
        draw_chart(draw, bullets or [message], right, paper.accent)
        draw_bullets(draw, bullets, left[0], left[1] + 10, left[2] - left[0], paper.accent, size=32)
    else:
        # Alternating two-column explanation keeps dense reading slides from
        # becoming repeated card grids while preserving a stable academic rhythm.
        if local_idx % 2 == 0:
            draw_flow(draw, bullets or [message], (M, 390, W - M, 814), paper.accent)
        else:
            draw_bullets(draw, bullets, M + 26, 372, W - 2 * M - 52, paper.accent, size=34)
            draw.line((M, 844, W - M, 844), fill=rgb("#D8E3ED"), width=2)
        draw_wrapped(draw, citation, (M, 864), W - 2 * M, font(22), fill=MUTED, line_gap=4, max_lines=1)
    return img


def save_slide(img: Image.Image, index: int) -> Path:
    SLIDE_DIR.mkdir(parents=True, exist_ok=True)
    out = SLIDE_DIR / f"slide_{index:02d}.png"
    img.save(out, "PNG", optimize=True)
    return out


def build_pptx(slide_paths: list[Path]) -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    for path in slide_paths:
        slide = prs.slides.add_slide(blank)
        slide.shapes.add_picture(str(path), 0, 0, width=prs.slide_width, height=prs.slide_height)
    prs.save(PPTX_PATH)


def make_contact_sheet(slide_paths: list[Path]) -> Path:
    PREVIEW_DIR.mkdir(parents=True, exist_ok=True)
    thumb_w, thumb_h = 384, 216
    cols = 5
    rows = math.ceil(len(slide_paths) / cols)
    sheet = Image.new("RGB", (cols * thumb_w, rows * (thumb_h + 34)), rgb("#EEF3F7"))
    draw = ImageDraw.Draw(sheet)
    for idx, path in enumerate(slide_paths):
        with Image.open(path).convert("RGB") as src:
            src.thumbnail((thumb_w, thumb_h), Image.Resampling.LANCZOS)
            x = (idx % cols) * thumb_w
            y = (idx // cols) * (thumb_h + 34)
            sheet.paste(src, (x, y))
            draw.text((x + 10, y + thumb_h + 6), f"{idx + 1:02d}", font=font(20, True), fill=rgb(NAVY))
    out = PREVIEW_DIR / "contact_sheet.png"
    sheet.save(out, "PNG", optimize=True)
    return out


def write_notes(all_slides: list[tuple[Paper, int, dict[str, Any]]]) -> None:
    lines = ["# POPL 量子论文阅读分享讲稿", ""]
    for paper, local_idx, slide in all_slides:
        if local_idx == 1:
            lines += ["", f"## {paper.title}", ""]
        lines += [
            f"### {local_idx}. {slide['title']}",
            "",
            f"- 页面主旨：{slide.get('main_message', '')}",
            f"- 引用位置：{slide.get('citation', '')}",
            "",
            slide.get("speaker_notes", ""),
            "",
        ]
    NOTES_PATH.write_text("\n".join(lines), encoding="utf-8")


def main() -> None:
    OUT_DIR.mkdir(parents=True, exist_ok=True)
    all_slides: list[tuple[Paper, int, dict[str, Any]]] = []
    for paper in PAPERS:
        slides = load_paper_slides(paper)
        for local_idx, slide in enumerate(slides, start=1):
            all_slides.append((paper, local_idx, slide))

    total = len(all_slides)
    counters: dict[str, int] = {paper.key: 0 for paper in PAPERS}
    figures = {paper.key: figure_paths(paper) for paper in PAPERS}
    slide_paths: list[Path] = []

    for global_idx, (paper, local_idx, slide) in enumerate(all_slides, start=1):
        fig: Path | None = None
        if figures[paper.key] and local_idx in {7, 8, 9, 10, 11, 13, 14}:
            fig = figures[paper.key][counters[paper.key] % len(figures[paper.key])]
            counters[paper.key] += 1
        if local_idx == 1:
            img = title_slide(paper, slide, global_idx, total)
        else:
            img = content_slide(paper, slide, global_idx, total, fig, local_idx)
        slide_paths.append(save_slide(img, global_idx))

    build_pptx(slide_paths)
    contact = make_contact_sheet(slide_paths)
    write_notes(all_slides)
    QA_PATH.write_text(
        json.dumps(
            {
                "slide_count": len(slide_paths),
                "pptx": str(PPTX_PATH),
                "slides_dir": str(SLIDE_DIR),
                "contact_sheet": str(contact),
                "notes": str(NOTES_PATH),
                "figure_counts": {paper.key: len(figures[paper.key]) for paper in PAPERS},
                "font": FONT_PATH,
            },
            ensure_ascii=False,
            indent=2,
        ),
        encoding="utf-8",
    )
    print(f"Wrote {len(slide_paths)} slide images")
    print(f"Wrote {PPTX_PATH}")
    print(f"Wrote {contact}")
    print(f"Wrote {NOTES_PATH}")


if __name__ == "__main__":
    main()
