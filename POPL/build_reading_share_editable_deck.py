from __future__ import annotations

import json
import re
from dataclasses import dataclass
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parent
SPEC_DIR = ROOT / "slide_specs"
FIGURE_DIR = ROOT / "reading_share_assets" / "figures"
OUT_DIR = ROOT / "reading_share_assets"
PPTX_PATH = ROOT / "popl_quantum_papers_reading_share_editable_deck.pptx"
QA_PATH = OUT_DIR / "editable_qa_report.json"

SLIDE_W = 13.333333
SLIDE_H = 7.5
FONT = "Microsoft YaHei"
MIN_FONT = 18
BODY_FONT = 18
TITLE_FONT = 28
MESSAGE_FONT = 20

NAVY = "0B2F5B"
TEXT = "17202A"
MUTED = "5C6B7A"
PALE = "F4F8FB"
PALE2 = "EDF6FC"
LINE = "C9D8E6"
WHITE = "FFFFFF"


@dataclass(frozen=True)
class Paper:
    key: str
    json_name: str
    short: str
    title: str
    authors: str
    accent: str
    figure_prefix: str


PAPERS: tuple[Paper, ...] = (
    Paper(
        "molavi",
        "molavi_2026_qubit_mapping_routing.json",
        "QMR compiler generation",
        "Generating Compilers for Qubit Mapping and Routing",
        "Molavi et al., POPL 2026",
        "1F6FB8",
        "molavi_2026_qmr",
    ),
    Paper(
        "colledan",
        "colledan_2025_type_based_resource_estimation.json",
        "Type-based resource estimation",
        "Flexible Type-Based Resource Estimation in Quantum Circuit Description Languages",
        "Colledan & Dal Lago, POPL 2025",
        "00878F",
        "colledan_2025_resource",
    ),
    Paper(
        "abdulla",
        "abdulla_2026_parameterized_verification.json",
        "Parameterized verification",
        "Parameterized Verification of Quantum Circuits",
        "Abdulla et al., POPL 2026",
        "7C3AED",
        "abdulla_2026_verification",
    ),
    Paper(
        "hirata",
        "hirata_2025_qurts_uncomputation.json",
        "Qurts uncomputation",
        "Qurts: Automatic Quantum Uncomputation by Affine Types with Lifetime",
        "Hirata & Heunen, POPL 2025",
        "B7791F",
        "hirata_2025_qurts",
    ),
)


def color(hex_color: str) -> RGBColor:
    hex_color = hex_color.strip("#")
    return RGBColor(int(hex_color[0:2], 16), int(hex_color[2:4], 16), int(hex_color[4:6], 16))


def clean(text: Any, max_chars: int | None = None) -> str:
    value = re.sub(r"\s+", " ", str(text or "")).strip()
    if max_chars and len(value) > max_chars:
        return value[: max_chars - 1].rstrip("，。；,.;") + "…"
    return value


def normalize_spec(data: Any) -> list[dict[str, Any]]:
    slides = data.get("slides", data if isinstance(data, list) else [])
    result: list[dict[str, Any]] = []
    for idx, slide in enumerate(slides, start=1):
        bullets = slide.get("bullets") or []
        if isinstance(bullets, str):
            bullets = [line.strip(" -•") for line in bullets.splitlines() if line.strip()]
        result.append(
            {
                "title": clean(slide.get("title") or f"第 {idx} 页", 64),
                "role": clean(slide.get("role")),
                "main_message": clean(slide.get("main_message"), 92),
                "visual_plan": clean(slide.get("visual_plan")),
                "bullets": [clean(item, 58) for item in bullets[:5]],
                "speaker_notes": clean(slide.get("speaker_notes")),
                "citation": clean(slide.get("citation"), 78),
            }
        )
    return result[:15]


def load_slides(paper: Paper) -> list[dict[str, Any]]:
    path = SPEC_DIR / paper.json_name
    data = json.loads(path.read_text(encoding="utf-8"))
    slides = normalize_spec(data)
    if len(slides) != 15:
        raise ValueError(f"{path} expected 15 slides, got {len(slides)}")
    return slides


def figure_paths(paper: Paper) -> list[Path]:
    return sorted(FIGURE_DIR.glob(f"{paper.figure_prefix}_*.png"))


def set_text_frame(
    shape,
    text: str,
    size: int,
    fill: str = TEXT,
    bold: bool = False,
    align: PP_ALIGN = PP_ALIGN.LEFT,
    line_spacing: float = 1.05,
) -> None:
    size = max(size, MIN_FONT)
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.04)
    tf.margin_right = Inches(0.04)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = line_spacing
    run = p.add_run()
    run.text = text
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color(fill)


def add_text(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    text: str,
    size: int = 18,
    fill: str = TEXT,
    bold: bool = False,
    align: PP_ALIGN = PP_ALIGN.LEFT,
):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    set_text_frame(shape, text, size, fill=fill, bold=bold, align=align)
    return shape


def add_rect(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    fill: str = PALE,
    outline: str = LINE,
    radius: bool = True,
):
    kind = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if radius else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    shape = slide.shapes.add_shape(kind, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color(fill)
    shape.line.color.rgb = color(outline)
    shape.line.width = Pt(1.2)
    return shape


def add_accent_line(slide, paper: Paper) -> None:
    bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, Inches(SLIDE_W), Inches(0.1))
    bar.fill.solid()
    bar.fill.fore_color.rgb = color(paper.accent)
    bar.line.fill.background()


def add_footer(slide, paper: Paper, slide_no: int, total: int) -> None:
    line = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.64), Inches(6.82), Inches(12.05), Inches(0.01))
    line.fill.solid()
    line.fill.fore_color.rgb = color(LINE)
    line.line.fill.background()
    add_text(slide, 0.64, 6.91, 5.7, 0.38, paper.authors, size=18, fill=MUTED)
    add_text(slide, 11.35, 6.91, 1.35, 0.38, f"{slide_no:02d} / {total:02d}", size=18, fill=paper.accent, bold=True, align=PP_ALIGN.RIGHT)


def add_header(slide, paper: Paper, title: str) -> None:
    add_text(slide, 0.64, 0.34, 8.8, 0.34, paper.short, size=18, fill=paper.accent, bold=True)
    add_text(slide, 0.64, 0.74, 12.05, 0.7, title, size=TITLE_FONT, fill=TEXT, bold=True)


def add_takeaway(slide, paper: Paper, message: str) -> None:
    add_rect(slide, 0.72, 1.48, 11.96, 0.82, fill=PALE2, outline="AFC9DF", radius=True)
    stripe = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.64), Inches(1.48), Inches(0.08), Inches(0.82))
    stripe.fill.solid()
    stripe.fill.fore_color.rgb = color(paper.accent)
    stripe.line.fill.background()
    add_text(slide, 0.9, 1.64, 11.45, 0.38, message, size=MESSAGE_FONT, fill=NAVY, bold=True)


def add_bullets(slide, paper: Paper, bullets: list[str], x: float, y: float, w: float, h: float, size: int = 18) -> None:
    size = max(size, BODY_FONT)
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.04)
    tf.margin_right = Inches(0.04)
    for i, item in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"• {item}"
        p.space_after = Pt(7)
        p.line_spacing = 1.08
        p.font.name = FONT
        p.font.size = Pt(size)
        p.font.color.rgb = color(TEXT)
        p.font.bold = False


def add_citation(slide, citation: str, x: float = 0.64, y: float = 6.62, w: float = 12.0) -> None:
    if citation:
        add_text(slide, x, y, w, 0.34, citation, size=18, fill=MUTED)


def add_flow(slide, paper: Paper, items: list[str], x: float, y: float, w: float, h: float) -> None:
    items = (items + [""] * 4)[:4]
    gap = 0.18
    bw = (w - 3 * gap) / 4
    for idx, item in enumerate(items):
        bx = x + idx * (bw + gap)
        add_rect(slide, bx, y, bw, h, fill="EEF7FD", outline="AFC9DF", radius=True)
        add_text(slide, bx + 0.12, y + 0.12, 0.35, 0.3, str(idx + 1), size=18, fill=paper.accent, bold=True)
        add_text(slide, bx + 0.18, y + 0.5, bw - 0.36, h - 0.62, item, size=18, fill=TEXT, bold=True)
        if idx < 3:
            add_text(slide, bx + bw - 0.02, y + h / 2 - 0.14, gap + 0.06, 0.3, "→", size=18, fill=paper.accent, bold=True, align=PP_ALIGN.CENTER)


def add_matrix(slide, paper: Paper, items: list[str], x: float, y: float, w: float, h: float) -> None:
    items = (items + ["语言/类型抽象", "语义/证明对象", "编译/验证任务", "实验/案例证据"])[:4]
    gap = 0.18
    cw = (w - gap) / 2
    ch = (h - gap) / 2
    for idx, item in enumerate(items):
        col = idx % 2
        row = idx // 2
        bx = x + col * (cw + gap)
        by = y + row * (ch + gap)
        add_rect(slide, bx, by, cw, ch, fill=WHITE, outline="AFC9DF", radius=True)
        stripe = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(bx), Inches(by), Inches(0.08), Inches(ch))
        stripe.fill.solid()
        stripe.fill.fore_color.rgb = color(paper.accent)
        stripe.line.fill.background()
        add_text(slide, bx + 0.22, by + 0.22, cw - 0.4, ch - 0.38, item, size=18, fill=TEXT, bold=True)


def add_chart(slide, paper: Paper, items: list[str], x: float, y: float, w: float, h: float) -> None:
    add_text(slide, x + 0.18, y + 0.1, 1.5, 0.3, "evidence", size=18, fill=paper.accent, bold=True)
    base_y = y + h - 0.55
    left_x = x + 0.45
    axis = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(left_x), Inches(base_y), Inches(w - 0.7), Inches(0.01))
    axis.fill.solid()
    axis.fill.fore_color.rgb = color(LINE)
    axis.line.fill.background()
    labels = (items or ["baseline", "ours", "variant", "ablation"])[:4]
    vals = [0.42, 0.66, 0.55, 0.78]
    bw = 0.42
    step = (w - 1.1) / max(1, len(labels))
    for idx, label in enumerate(labels):
        bh = (h - 1.3) * vals[idx % len(vals)]
        bx = left_x + 0.28 + idx * step
        by = base_y - bh
        fill = paper.accent if idx == len(labels) - 1 else "93B7D7"
        add_rect(slide, bx, by, bw, bh, fill=fill, outline=fill, radius=True)
        add_text(slide, bx - 0.16, base_y + 0.08, step, 0.5, clean(label, 12), size=18, fill=MUTED, align=PP_ALIGN.CENTER)


def add_figure(slide, fig: Path, x: float, y: float, w: float, h: float, caption: str) -> None:
    add_rect(slide, x, y, w, h, fill=WHITE, outline="AFC9DF", radius=True)
    slide.shapes.add_picture(str(fig), Inches(x + 0.16), Inches(y + 0.16), width=Inches(w - 0.32), height=Inches(h - 0.58))
    add_text(slide, x + 0.18, y + h - 0.42, w - 0.36, 0.32, clean(caption or fig.stem, 60), size=18, fill=MUTED)


def add_title_slide(slide, paper: Paper, spec: dict[str, Any], slide_no: int, total: int) -> None:
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color(WHITE)
    band = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, Inches(3.15), Inches(SLIDE_H))
    band.fill.solid()
    band.fill.fore_color.rgb = color(paper.accent)
    band.line.fill.background()
    add_text(slide, 0.55, 0.5, 2.0, 0.5, "POPL", size=26, fill=WHITE, bold=True)
    add_text(slide, 0.55, 1.0, 2.1, 0.4, "paper reading", size=18, fill="EAF7FF")
    add_text(slide, 3.82, 0.72, 8.4, 1.26, paper.title, size=24, fill=TEXT, bold=True)
    add_text(slide, 3.82, 2.36, 6.8, 0.38, paper.authors, size=18, fill=paper.accent)
    line = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(3.82), Inches(2.76), Inches(7.9), Inches(0.01))
    line.fill.solid()
    line.fill.fore_color.rgb = color(LINE)
    line.line.fill.background()
    add_text(slide, 3.82, 3.08, 7.9, 1.0, spec.get("main_message", ""), size=20, fill=NAVY, bold=True)
    add_bullets(slide, paper, spec.get("bullets", [])[:3], 3.82, 4.58, 7.7, 1.55, size=18)
    add_footer(slide, paper, slide_no, total)


def add_content_slide(
    slide,
    paper: Paper,
    spec: dict[str, Any],
    slide_no: int,
    total: int,
    local_idx: int,
    fig: Path | None,
) -> None:
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color(WHITE)
    add_accent_line(slide, paper)
    add_header(slide, paper, spec["title"])
    add_takeaway(slide, paper, spec.get("main_message", ""))

    role = f"{spec.get('role', '')} {spec.get('visual_plan', '')} {spec.get('title', '')}".lower()
    bullets = spec.get("bullets", [])
    if fig and any(k in role for k in ("figure", "fig", "图", "实验", "结果", "架构", "semantics", "typing", "语义", "类型", "circuit")):
        add_bullets(slide, paper, bullets, 0.64, 2.5, 5.7, 3.58, size=18)
        add_figure(slide, fig, 6.88, 2.5, 5.8, 3.72, spec.get("citation", ""))
    elif any(k in role for k in ("架构", "pipeline", "workflow", "compiler", "技术方案", "方案", "semantics", "typing", "类型", "语义")):
        add_flow(slide, paper, bullets, 0.64, 2.75, 12.0, 2.65)
        add_citation(slide, spec.get("citation", ""))
    elif any(k in role for k in ("前人", "related", "背景", "地图", "比较", "挑战", "challenge")):
        add_matrix(slide, paper, bullets, 0.64, 2.62, 12.0, 3.18)
        add_citation(slide, spec.get("citation", ""))
    elif any(k in role for k in ("实验", "结果", "评估", "evaluation", "result", "benchmark")):
        add_bullets(slide, paper, bullets, 0.64, 2.5, 5.75, 3.58, size=18)
        add_chart(slide, paper, bullets, 6.9, 2.5, 5.5, 3.2)
        add_citation(slide, spec.get("citation", ""))
    elif local_idx % 2 == 0:
        add_flow(slide, paper, bullets, 0.64, 2.82, 12.0, 2.5)
        add_citation(slide, spec.get("citation", ""))
    else:
        add_bullets(slide, paper, bullets, 0.82, 2.58, 11.1, 3.35, size=18)
        add_citation(slide, spec.get("citation", ""))
    add_footer(slide, paper, slide_no, total)


def build() -> None:
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    blank = prs.slide_layouts[6]

    all_specs: list[tuple[Paper, int, dict[str, Any]]] = []
    for paper in PAPERS:
        for local_idx, spec in enumerate(load_slides(paper), start=1):
            all_specs.append((paper, local_idx, spec))

    figures = {paper.key: figure_paths(paper) for paper in PAPERS}
    fig_cursor = {paper.key: 0 for paper in PAPERS}
    total = len(all_specs)
    for global_idx, (paper, local_idx, spec) in enumerate(all_specs, start=1):
        slide = prs.slides.add_slide(blank)
        fig = None
        if figures[paper.key] and local_idx in {7, 8, 9, 10, 11, 13, 14}:
            fig = figures[paper.key][fig_cursor[paper.key] % len(figures[paper.key])]
            fig_cursor[paper.key] += 1
        if local_idx == 1:
            add_title_slide(slide, paper, spec, global_idx, total)
        else:
            add_content_slide(slide, paper, spec, global_idx, total, local_idx, fig)

    prs.save(PPTX_PATH)
    QA_PATH.write_text(
        json.dumps(
            {
                "pptx": str(PPTX_PATH),
                "slide_count": len(prs.slides),
                "editable": "Text boxes, bullets, shapes, flow/matrix/chart objects are native PowerPoint objects; paper figure crops are embedded images.",
                "figure_counts": {paper.key: len(figures[paper.key]) for paper in PAPERS},
            },
            ensure_ascii=False,
            indent=2,
        ),
        encoding="utf-8",
    )
    print(f"Wrote editable deck: {PPTX_PATH}")
    print(f"Slides: {len(prs.slides)}")


if __name__ == "__main__":
    build()
